"""
文件解析服务 - 支持12种格式16种扩展名
PDF/Word/Excel/PPT/CSV/TXT/RTF/HTML/Markdown/XMind + OCR扫描件
"""

import os
import logging
from typing import Optional

logger = logging.getLogger(__name__)

# 支持的文件格式
SUPPORTED_FORMATS = {
    '.pdf': 'application/pdf',
    '.docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
    '.doc': 'application/msword',
    '.xlsx': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
    '.xls': 'application/vnd.ms-excel',
    '.pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
    '.ppt': 'application/vnd.ms-powerpoint',
    '.csv': 'text/csv',
    '.txt': 'text/plain',
    '.rtf': 'application/rtf',
    '.html': 'text/html',
    '.htm': 'text/html',
    '.md': 'text/markdown',
    '.markdown': 'text/markdown',
    '.xmind': 'application/xmind',
    '.png': 'image/png',
    '.jpg': 'image/jpeg',
    '.jpeg': 'image/jpeg',
    '.bmp': 'image/bmp',
    '.tiff': 'image/tiff',
    '.tif': 'image/tiff',
    '.gif': 'image/gif',
}


def get_supported_extensions() -> list:
    return list(SUPPORTED_FORMATS.keys())


def is_supported_file(filename: str) -> bool:
    ext = os.path.splitext(filename)[1].lower()
    return ext in SUPPORTED_FORMATS


async def extract_text_from_file(file_path: str, max_length: int = 15000) -> Optional[str]:
    """
    从文件中提取文本内容
    支持: PDF/Word/Excel/PPT/CSV/TXT/RTF/HTML/Markdown/XMind/图片OCR
    """
    if not os.path.exists(file_path):
        logger.error(f"文件不存在: {file_path}")
        return None

    ext = os.path.splitext(file_path)[1].lower()

    try:
        if ext == '.txt':
            return _read_text_file(file_path, max_length)
        elif ext == '.rtf':
            return _extract_rtf(file_path, max_length)
        elif ext in ('.html', '.htm'):
            return _extract_html(file_path, max_length)
        elif ext in ('.md', '.markdown'):
            return _read_text_file(file_path, max_length)
        elif ext == '.csv':
            return _extract_csv(file_path, max_length)
        elif ext == '.pdf':
            return _extract_pdf(file_path, max_length)
        elif ext in ('.doc', '.docx'):
            return _extract_word(file_path, max_length)
        elif ext in ('.xlsx', '.xls'):
            return _extract_excel(file_path, max_length)
        elif ext in ('.pptx', '.ppt'):
            return _extract_ppt(file_path, max_length)
        elif ext == '.xmind':
            return _extract_xmind(file_path, max_length)
        elif ext in ('.png', '.jpg', '.jpeg', '.bmp', '.tiff', '.tif', '.gif'):
            return _extract_image_ocr(file_path, max_length)
        else:
            logger.warning(f"不支持的文件格式: {ext}")
            return None
    except Exception as e:
        logger.error(f"文件解析失败 [{ext}]: {e}", exc_info=True)
        return None


def _read_text_file(file_path: str, max_length: int) -> str:
    """读取纯文本文件，自动探测编码"""
    encodings = ['utf-8', 'gbk', 'gb2312', 'gb18030', 'latin-1']
    for enc in encodings:
        try:
            with open(file_path, 'r', encoding=enc) as f:
                text = f.read(max_length)
            return text
        except (UnicodeDecodeError, UnicodeError):
            continue
    raise ValueError(f"无法解码文件: {file_path}")


def _extract_rtf(file_path: str, max_length: int) -> str:
    """提取RTF文件内容"""
    try:
        from striprtf.striprtf import rtf_to_text
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            rtf_content = f.read()
        text = rtf_to_text(rtf_content)
        return text[:max_length]
    except ImportError:
        logger.warning("striprtf未安装，尝试纯文本读取")
        return _read_text_file(file_path, max_length)


def _extract_html(file_path: str, max_length: int) -> str:
    """提取HTML文件内容"""
    try:
        from bs4 import BeautifulSoup
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            html = f.read()
        soup = BeautifulSoup(html, 'html.parser')
        # 移除script和style
        for tag in soup(['script', 'style']):
            tag.decompose()
        text = soup.get_text(separator='\n', strip=True)
        return text[:max_length]
    except ImportError:
        # 回退：简单正则去标签
        import re
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            html = f.read()
        text = re.sub(r'<[^>]+>', ' ', html)
        text = re.sub(r'\s+', ' ', text).strip()
        return text[:max_length]


def _extract_csv(file_path: str, max_length: int) -> str:
    """提取CSV内容为可读文本"""
    import csv
    encodings = ['utf-8', 'gbk', 'gb2312', 'gb18030', 'latin-1']
    for enc in encodings:
        try:
            with open(file_path, 'r', encoding=enc) as f:
                reader = csv.reader(f)
                lines = []
                for i, row in enumerate(reader):
                    if i >= 200:  # 限制行数
                        break
                    lines.append(' | '.join(row))
                text = '\n'.join(lines)
                return text[:max_length]
        except (UnicodeDecodeError, UnicodeError):
            continue
    raise ValueError(f"无法解码CSV文件: {file_path}")


def _extract_pdf(file_path: str, max_length: int) -> str:
    """提取PDF内容，扫描件自动降级OCR"""
    try:
        import fitz  # PyMuPDF
        doc = fitz.open(file_path)
        text = ""
        has_text = False
        for page in doc:
            page_text = page.get_text()
            if page_text.strip():
                has_text = True
            text += page_text
            if len(text) > max_length:
                break
        doc.close()

        if has_text and text.strip():
            return text[:max_length]

        # 文本层为空 → 扫描件 PDF，降级 OCR
        logger.info(f"PDF无文本层，使用OCR识别: {file_path}")
        return _ocr_pdf(file_path, max_length)
    except ImportError:
        logger.warning("PyMuPDF未安装，无法解析PDF")
        return None


def _extract_word(file_path: str, max_length: int) -> str:
    """提取Word文档内容 (.docx / .doc)"""
    ext = os.path.splitext(file_path)[1].lower()

    # .doc 旧格式: 用 textutil (macOS) 或 antiword (Linux) 转换
    if ext == '.doc':
        import subprocess
        import tempfile
        try:
            # macOS textutil
            result = subprocess.run(
                ['textutil', '-convert', 'txt', '-stdout', file_path],
                capture_output=True, text=True, timeout=30
            )
            if result.returncode == 0 and result.stdout.strip():
                return result.stdout[:max_length]
        except (FileNotFoundError, subprocess.TimeoutExpired):
            pass
        try:
            # Linux antiword
            result = subprocess.run(
                ['antiword', file_path],
                capture_output=True, text=True, timeout=30
            )
            if result.returncode == 0 and result.stdout.strip():
                return result.stdout[:max_length]
        except (FileNotFoundError, subprocess.TimeoutExpired):
            pass
        try:
            # Linux catdoc
            result = subprocess.run(
                ['catdoc', file_path],
                capture_output=True, text=True, timeout=30
            )
            if result.returncode == 0 and result.stdout.strip():
                return result.stdout[:max_length]
        except (FileNotFoundError, subprocess.TimeoutExpired):
            pass
        logger.error(f"无法解析 .doc 文件 (需要 textutil/antiword/catdoc): {file_path}")
        return None

    # .docx 新格式
    try:
        import docx
        doc = docx.Document(file_path)
        text = "\n".join([p.text for p in doc.paragraphs if p.text.strip()])
        return text[:max_length]
    except ImportError:
        logger.warning("python-docx未安装，无法解析Word文档")
        return None


def _extract_excel(file_path: str, max_length: int) -> str:
    """提取Excel内容为可读文本"""
    try:
        import openpyxl
        wb = openpyxl.load_workbook(file_path, data_only=True)
        lines = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            lines.append(f"=== 工作表: {sheet_name} ===")
            for row in ws.iter_rows(max_row=200, values_only=True):
                row_text = ' | '.join([str(c) if c is not None else '' for c in row])
                if row_text.strip(' |'):
                    lines.append(row_text)
                if len('\n'.join(lines)) > max_length:
                    break
        wb.close()
        return '\n'.join(lines)[:max_length]
    except ImportError:
        logger.warning("openpyxl未安装，尝试xlrd")
        try:
            import xlrd
            wb = xlrd.open_workbook(file_path)
            lines = []
            for sheet in wb.sheets():
                lines.append(f"=== 工作表: {sheet.name} ===")
                for row_idx in range(min(sheet.nrows, 200)):
                    row = [str(sheet.cell_value(row_idx, col_idx)) for col_idx in range(sheet.ncols)]
                    lines.append(' | '.join(row))
            return '\n'.join(lines)[:max_length]
        except ImportError:
            logger.warning("openpyxl和xlrd均未安装，无法解析Excel")
            return None


def _extract_ppt(file_path: str, max_length: int) -> str:
    """提取PPT内容"""
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        lines = []
        for i, slide in enumerate(prs.slides):
            lines.append(f"--- 幻灯片 {i + 1} ---")
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text.strip():
                    lines.append(shape.text)
        return '\n'.join(lines)[:max_length]
    except ImportError:
        logger.warning("python-pptx未安装，无法解析PPT")
        return None


def _extract_xmind(file_path: str, max_length: int) -> str:
    """提取XMind思维导图内容"""
    try:
        import zipfile
        import json
        with zipfile.ZipFile(file_path, 'r') as z:
            # XMind 8+ 使用JSON格式
            for name in z.namelist():
                if name.endswith('content.json'):
                    data = json.loads(z.read(name))
                    lines = _parse_xmind_json(data)
                    return '\n'.join(lines)[:max_length]
            # XMind旧格式
            for name in z.namelist():
                if 'content.xml' in name:
                    import xml.etree.ElementTree as ET
                    xml_data = z.read(name)
                    root = ET.fromstring(xml_data)
                    lines = []
                    _parse_xmind_xml(root, lines, 0)
                    return '\n'.join(lines)[:max_length]
        return None
    except Exception as e:
        logger.warning(f"XMind解析失败: {e}")
        return None


def _parse_xmind_json(data, depth=0) -> list:
    """解析XMind JSON结构"""
    lines = []
    if isinstance(data, list):
        for item in data:
            lines.extend(_parse_xmind_json(item, depth))
    elif isinstance(data, dict):
        title = data.get('title', '')
        if title:
            prefix = '  ' * depth
            lines.append(f"{prefix}- {title}")
        children = data.get('children', {}).get('attached', [])
        for child in children:
            lines.extend(_parse_xmind_json(child, depth + 1))
    return lines


def _parse_xmind_xml(element, lines, depth):
    """解析XMind XML结构"""
    tag = element.tag.split('}')[-1] if '}' in element.tag else element.tag
    if tag == 'topic':
        title_el = element.find('{http://www.xmind.net/xmind}title')
        if title_el is None:
            title_el = element.find('title')
        if title_el is not None and title_el.text:
            prefix = '  ' * depth
            lines.append(f"{prefix}- {title_el.text}")
    for child in element:
        _parse_xmind_xml(child, lines, depth + 1)


def _group_ocr_items(items: list, y_threshold: float = 30.0) -> list:
    """将OCR结果按Y坐标聚类成行，每行按X排序"""
    if not items:
        return []
    items_sorted = sorted(items, key=lambda a: a['y_center'])
    rows = []
    current_row = [items_sorted[0]]
    for item in items_sorted[1:]:
        if abs(item['y_center'] - current_row[0]['y_center']) < y_threshold:
            current_row.append(item)
        else:
            rows.append(current_row)
            current_row = [item]
    rows.append(current_row)
    for row in rows:
        row.sort(key=lambda a: a['x1'])
    return rows


def _detect_table_columns(rows: list, min_cols: int = 3, min_rows: int = 3) -> tuple:
    """检测表格区域：返回 (table_start_idx, table_end_idx, col_boundaries) 或 None
    col_boundaries 是表头每列的 (x1, x2)
    """
    # 找多列行（>= min_cols）
    multi_col_rows = []
    for i, row in enumerate(rows):
        if len(row) >= min_cols:
            multi_col_rows.append(i)
    if len(multi_col_rows) < min_rows:
        return None
    # 找最长的连续段（允许中间有1行间隔）
    best_start, best_end = multi_col_rows[0], multi_col_rows[0]
    cur_start, cur_end = multi_col_rows[0], multi_col_rows[0]
    for i in range(1, len(multi_col_rows)):
        if multi_col_rows[i] - multi_col_rows[i - 1] <= 3:
            cur_end = multi_col_rows[i]
        else:
            if cur_end - cur_start > best_end - best_start:
                best_start, best_end = cur_start, cur_end
            cur_start, cur_end = multi_col_rows[i], multi_col_rows[i]
    if cur_end - cur_start > best_end - best_start:
        best_start, best_end = cur_start, cur_end
    if best_end - best_start + 1 < min_rows:
        return None
    # 用众数（出现最频繁的列数）作为表格列数，避免OCR多分item导致列数偏大
    from collections import Counter
    table_rows = rows[best_start:best_end + 1]
    col_counts = Counter(len(r) for r in table_rows)
    max_cols = col_counts.most_common(1)[0][0]
    if max_cols < min_cols:
        max_cols = max(len(r) for r in table_rows)
    # 表头行：列数==max_cols 且 Y坐标最小（最靠上的多列行）
    header_candidates = [(i, r) for i, r in enumerate(table_rows) if len(r) == max_cols]
    header_idx_in_table = min(header_candidates, key=lambda x: x[1][0]['y_center'])[0] if header_candidates else 0
    header_row = table_rows[header_idx_in_table]
    # 表头必须是表格区域的第一行：调整 best_start 到表头位置
    best_start = best_start + header_idx_in_table
    header_idx_in_table = 0
    col_boundaries = [(item['x1'], item['x2']) for item in header_row]
    return (best_start, best_end, col_boundaries, header_idx_in_table)


def _assign_to_columns(row_items: list, col_boundaries: list) -> list:
    """将一行的OCR items分配到对应列，用左边界(x1)匹配列"""
    num_cols = len(col_boundaries)
    cells = [''] * num_cols
    # 列分界点：相邻列边界的中点
    col_dividers = []
    for ci in range(num_cols - 1):
        col_dividers.append((col_boundaries[ci][1] + col_boundaries[ci + 1][0]) / 2)
    for item in row_items:
        best_col = 0
        for ci, divider in enumerate(col_dividers):
            if item['x1'] > divider:
                best_col = ci + 1
            else:
                break
        if cells[best_col]:
            cells[best_col] += ' ' + item['text']
        else:
            cells[best_col] = item['text']
    return cells


def _merge_continuation_rows(table_rows: list, col_boundaries: list, min_cols: int = 3) -> list:
    """合并续行到上一行：列数 < min_cols 的行合并到上一行对应列
    table_rows[0] 是表头，不参与合并
    """
    if len(table_rows) <= 1:
        return table_rows
    merged = [table_rows[0]]  # 表头保持
    for row in table_rows[1:]:
        if len(row) < min_cols:
            # 续行：合并到上一行
            cells = _assign_to_columns(row, col_boundaries)
            if isinstance(merged[-1], dict) and merged[-1].get('_merged'):
                prev_cells = merged[-1]['_cells']
            else:
                prev_cells = _assign_to_columns(merged[-1], col_boundaries)
            for ci in range(len(cells)):
                if cells[ci]:
                    if prev_cells[ci]:
                        prev_cells[ci] += ' ' + cells[ci]
                    else:
                        prev_cells[ci] = cells[ci]
            merged[-1] = {'_cells': prev_cells, '_merged': True}
        else:
            merged.append(row)
    return merged


def _build_markdown_table(table_rows: list, col_boundaries: list) -> list:
    """将表格行构建为Markdown表格行列表"""
    num_cols = len(col_boundaries)
    lines = []
    header_done = False
    for row in table_rows:
        if isinstance(row, dict) and row.get('_merged'):
            cells = row['_cells']
        else:
            cells = _assign_to_columns(row, col_boundaries)
        lines.append('| ' + ' | '.join(cells) + ' |')
        if not header_done:
            lines.append('| ' + ' | '.join(['---'] * num_cols) + ' |')
            header_done = True
    return lines


def _ocr_pdf(file_path: str, max_length: int) -> str:
    """对扫描件PDF做OCR：逐页转图片 → PaddleOCR识别 → 重建表格"""
    import tempfile
    import os
    try:
        import fitz
    except ImportError:
        logger.warning("PyMuPDF未安装，无法OCR PDF")
        return None

    try:
        from paddleocr import PaddleOCR
        ocr = PaddleOCR(lang='ch', use_doc_orientation_classify=False, use_doc_unwarping=False)
    except Exception as e:
        logger.warning(f"PaddleOCR初始化失败: {e}")
        return _ocr_pdf_fallback(file_path, max_length)

    doc = fitz.open(file_path)
    all_lines = []
    for page_idx in range(len(doc)):
        page = doc[page_idx]
        mat = fitz.Matrix(3.0, 3.0)
        pix = page.get_pixmap(matrix=mat)
        tmp = tempfile.NamedTemporaryFile(suffix='.png', delete=False)
        pix.save(tmp.name)

        try:
            # 获取OCR结果（含坐标）
            items = []
            if hasattr(ocr, 'predict'):
                results = ocr.predict(tmp.name)
                for r in results:
                    res = r.json if hasattr(r, 'json') else r
                    if isinstance(res, dict) and 'res' in res:
                        res = res['res']
                    dt_polys = res.get('dt_polys', []) if isinstance(res, dict) else []
                    rec_texts = res.get('rec_texts', []) if isinstance(res, dict) else []
                    rec_scores = res.get('rec_scores', []) if isinstance(res, dict) else []
                    for poly, text, score in zip(dt_polys, rec_texts, rec_scores):
                        xs = [p[0] for p in poly]
                        ys = [p[1] for p in poly]
                        items.append({
                            'text': text,
                            'x1': min(xs), 'x2': max(xs),
                            'y1': min(ys), 'y2': max(ys),
                            'y_center': (min(ys) + max(ys)) / 2,
                            'x_center': (min(xs) + max(xs)) / 2,
                        })
            else:
                result = ocr.ocr(tmp.name)
                if result and result[0]:
                    for line in result[0]:
                        text = line[1][0] if isinstance(line[1], (list, tuple)) else str(line[1])
                        poly = line[0]
                        xs = [p[0] for p in poly]
                        ys = [p[1] for p in poly]
                        items.append({
                            'text': text,
                            'x1': min(xs), 'x2': max(xs),
                            'y1': min(ys), 'y2': max(ys),
                            'y_center': (min(ys) + max(ys)) / 2,
                            'x_center': (min(xs) + max(xs)) / 2,
                        })

            if not items:
                continue

            all_lines.append(f"--- 第{page_idx + 1}页 ---")

            # 按Y坐标聚类成行
            rows = _group_ocr_items(items)

            # 检测表格区域
            table_info = _detect_table_columns(rows)
            if table_info:
                t_start, t_end, col_boundaries, header_idx_in_table = table_info
                # 表格前的非表格行
                for ri in range(t_start):
                    row_texts = [item['text'] for item in rows[ri]]
                    all_lines.append(' '.join(row_texts))
                # 表格行
                table_rows = rows[t_start:t_end + 1]
                # 表头前的行作为普通文本
                for i in range(header_idx_in_table):
                    all_lines.append(' '.join(item['text'] for item in table_rows[i]))
                # 表头 + 表头后的行
                final_rows = [table_rows[header_idx_in_table]] + table_rows[header_idx_in_table + 1:]
                # 合并续行
                final_rows = _merge_continuation_rows(final_rows, col_boundaries)
                md_lines = _build_markdown_table(final_rows, col_boundaries)
                all_lines.extend(md_lines)
                # 表格后的非表格行
                for ri in range(t_end + 1, len(rows)):
                    row_texts = [item['text'] for item in rows[ri]]
                    all_lines.append(' '.join(row_texts))
            else:
                # 没有表格，直接按行输出
                for row in rows:
                    row_texts = [item['text'] for item in row]
                    all_lines.append(' '.join(row_texts))

        except Exception as e:
            logger.warning(f"PDF第{page_idx+1}页OCR失败: {e}")
        finally:
            os.unlink(tmp.name)

    try:
        doc.close()
    except Exception:
        pass

    full_text = '\n'.join(all_lines)
    if not full_text.strip():
        return _ocr_pdf_fallback(file_path, max_length)
    return full_text[:max_length]


def _ocr_pdf_fallback(file_path: str, max_length: int) -> str:
    """PDF OCR降级：EasyOCR → Tesseract"""
    import tempfile
    import os
    try:
        import fitz
    except ImportError:
        return None

    # EasyOCR
    try:
        import easyocr
        reader = easyocr.Reader(['ch_sim', 'en'], gpu=False, verbose=False)
        doc = fitz.open(file_path)
        all_lines = []
        for page_idx in range(len(doc)):
            page = doc[page_idx]
            mat = fitz.Matrix(3.0, 3.0)
            pix = page.get_pixmap(matrix=mat)
            tmp = tempfile.NamedTemporaryFile(suffix='.png', delete=False)
            pix.save(tmp.name)
            result = reader.readtext(tmp.name, detail=0)
            if result:
                all_lines.append(f"--- 第{page_idx + 1}页 ---")
                all_lines.extend(result)
            os.unlink(tmp.name)
        doc.close()
        full_text = '\n'.join(all_lines)
        if full_text.strip():
            return full_text[:max_length]
    except Exception as e:
        logger.warning(f"EasyOCR PDF降级失败: {e}")

    # Tesseract
    try:
        import pytesseract
        from PIL import Image
        import io
        doc = fitz.open(file_path)
        all_lines = []
        for page_idx in range(len(doc)):
            page = doc[page_idx]
            mat = fitz.Matrix(3.0, 3.0)
            pix = page.get_pixmap(matrix=mat)
            img = Image.open(io.BytesIO(pix.tobytes('png')))
            text = pytesseract.image_to_string(img, lang='chi_sim+eng')
            if text.strip():
                all_lines.append(f"--- 第{page_idx + 1}页 ---")
                all_lines.append(text)
        doc.close()
        full_text = '\n'.join(all_lines)
        return full_text[:max_length] if full_text.strip() else None
    except Exception as e:
        logger.error(f"所有OCR引擎均失败: {e}")
        return None


def _extract_image_ocr(file_path: str, max_length: int) -> str:
    """OCR图片文字识别（中英文）- 首选PaddleOCR, 降级EasyOCR, 最后Tesseract"""
    # 1. PaddleOCR (中文准确率最高)
    try:
        from paddleocr import PaddleOCR
        ocr = PaddleOCR(lang='ch', use_doc_orientation_classify=False, use_doc_unwarping=False)

        # PaddleOCR 3.x 用 predict，2.x 用 ocr
        if hasattr(ocr, 'predict'):
            results = ocr.predict(file_path)
            lines = []
            for r in results:
                res = r.json if hasattr(r, 'json') else r
                if isinstance(res, dict) and 'res' in res:
                    res = res['res']
                rec_texts = res.get('rec_texts', []) if isinstance(res, dict) else []
                lines.extend(rec_texts)
        else:
            result = ocr.ocr(file_path)
            lines = []
            if result and result[0]:
                for line in result[0]:
                    text = line[1][0] if isinstance(line[1], (list, tuple)) else str(line[1])
                    lines.append(text)

        full_text = '\n'.join(lines)
        if full_text.strip():
            return full_text[:max_length]
    except Exception as e:
        logger.warning(f"PaddleOCR识别失败: {e}, 尝试EasyOCR")

    # 2. EasyOCR (降级)
    try:
        import easyocr
        reader = easyocr.Reader(['ch_sim', 'en'], gpu=False, verbose=False)
        result = reader.readtext(file_path, detail=0)
        full_text = '\n'.join(result)
        return full_text[:max_length]
    except ImportError:
        logger.warning("EasyOCR未安装，尝试pytesseract")
    except Exception as e:
        logger.warning(f"EasyOCR识别失败: {e}, 尝试pytesseract")

    # 3. Tesseract (最后兜底)
    try:
        import pytesseract
        from PIL import Image
        img = Image.open(file_path)
        text = pytesseract.image_to_string(img, lang='chi_sim+eng')
        return text[:max_length]
    except Exception as e:
        logger.error(f"所有OCR引擎均失败: {e}")
        return None
